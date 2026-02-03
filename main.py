
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 1. Create the Presentation Object
    prs = Presentation()

    # Helper function to add a slide with title and bullet points
    def add_slide(prs, title, content_list, notes=""):
        # Layout 1 is "Title and Content"
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set Content (Bullet points)
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            
            # check for sub-bullets (simple heuristic: starts with "-")
            if item.startswith("    -"):
                p.level = 1
                p.text = item.replace("    -", "").strip()

        # Add notes if any
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # --- SLIDE 1: TITLE ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "House Price Prediction using Linear Regression"
    slide.placeholders[1].text = "A Complete ML Pipeline with FastAPI Deployment\n\nHaroon Khan\nBS Computer Science (Part-III)\nUniversity of Sindh Laar Campus, Badin"

    # --- SLIDE 2: PROJECT GOAL ---
    content_2 = [
        "Objective: Predict continuous selling prices of houses.",
        "Key Features Used:",
        "    - Overall Quality & Ground Living Area",
        "    - Garage Capacity, Basement Area, Year Built",
        "Deliverable: Realistic price estimates for unseen houses.",
        "Scope: Full ML Pipeline:",
        "    - Data → Model → Evaluation → Deployment (FastAPI)",
        "[PLACEHOLDER: Add Pipeline Diagram Image Here]"
    ]
    add_slide(prs, "Project Goal", content_2)

    # --- SLIDE 3: WHY LINEAR REGRESSION ---
    content_3 = [
        "Model Choice: Linear Regression",
        "Reasoning:",
        "    - Target is continuous numeric (Price).",
        "    - Assumes linear relationship between features and price.",
        "Advantages:",
        "    - Highly interpretable (coefficients show feature importance).",
        "    - Efficient for small/medium tabular datasets.",
        "    - Perfect baseline before trying complex models (e.g., XGBoost).",
        "[PLACEHOLDER: Add Scatter Plot Image Here]"
    ]
    add_slide(prs, "Why Linear Regression?", content_3)

    # --- SLIDE 4: MATH & METRICS (UPDATED) ---
    content_4 = [
        "The Prediction Formula:",
        "    - y_pred = w0 + w1*x1 + w2*x2 ... + wn*xn",
        "    - Goal: Minimize error between y_pred and y_actual.",
        "Evaluation Metrics (Crucial):",
        "    - RMSE (Root Mean Squared Error):",
        "         * Measures average error in dollars ($).",
        "         * Lower is better.",
        "    - R² (R-Squared):",
        "         * Measures 'Goodness of Fit' (0 to 1).",
        "         * 0.83 means we explain 83% of price variation.",
        "[PLACEHOLDER: Add Residuals Graph Image Here]"
    ]
    add_slide(prs, "How It Works: Math & Metrics", content_4)

    # --- SLIDE 5: PIPELINE OVERVIEW ---
    content_5 = [
        "1. Data Loading:",
        "    - Ames Housing Dataset (Kaggle).",
        "2. EDA (Exploratory Data Analysis):",
        "    - Analyzed distribution of 'SalePrice'.",
        "    - Checked correlations to find best predictors.",
        "3. Preprocessing:",
        "    - Handling Missing Values (Imputation).",
        "    - Feature Scaling (StandardScaler) - Essential for Regression.",
        "[PLACEHOLDER: Add Correlation Heatmap Image Here]"
    ]
    add_slide(prs, "ML Pipeline: Data to Features", content_5)

    # --- SLIDE 6: TRAINING & RESULTS ---
    content_6 = [
        "Data Splitting:",
        "    - 70% Train, 15% Validation, 15% Test.",
        "Training:",
        "    - Model trained on Scaled Data.",
        "Results:",
        "    - Validation R²: ~0.83 (Good fit).",
        "    - Test R²: ~0.77 (Good generalization).",
        "    - RMSE provides the error margin in actual currency.",
        "[PLACEHOLDER: Add Screenshot of Accuracy Score Output]"
    ]
    add_slide(prs, "Model Training & Evaluation", content_6)

    # --- SLIDE 7: DEPLOYMENT ---
    content_7 = [
        "Persistence:",
        "    - Model & Scaler saved using joblib.",
        "Backend (FastAPI):",
        "    - API endpoint accepts raw features.",
        "    - Applies scaling -> Predicts price -> Returns JSON.",
        "Frontend:",
        "    - Simple HTML form for user input.",
        "[PLACEHOLDER: Add Screenshot of FastAPI Swagger UI or Form]"
    ]
    add_slide(prs, "Deployment: FastAPI & Web UI", content_7)

    # --- SLIDE 8: SUMMARY ---
    content_8 = [
        "Summary:",
        "1. Linear Regression is effective for continuous price prediction.",
        "2. Feature Scaling is mandatory for this algorithm.",
        "3. Metrics: R² for fit, RMSE for dollar-error.",
        "4. Successfully demonstrated Data -> Model -> Web API workflow.",
        "Questions?"
    ]
    add_slide(prs, "Key Takeaways", content_8)

    # Save the file
    file_name = "House_Price_Prediction.pptx"
    prs.save(file_name)
    print(f"Successfully created {file_name}")

if __name__ == "__main__":
    create_presentation()
